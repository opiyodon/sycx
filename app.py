from flask import Flask, request, jsonify
from summarizer import summarize_documents, enrich_summary_with_visuals
from model import get_model
import os
from dotenv import load_dotenv
import base64
import io
from pdfminer.high_level import extract_text
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation
from PIL import Image
import pytesseract

load_dotenv()
app = Flask(__name__)
summarization_model = get_model()

@app.route('/summarize', methods=['POST'])
def summarize():
    data = request.json
    user_id = data.get('user_id')
    documents = data.get('documents')
    merge_summaries = data.get('merge_summaries', False)
    summary_depth = data.get('summary_depth', 0.5)
    language = data.get('language', 'en')
    enable_visual_enhancements = data.get('enable_visual_enhancements', True)

    if not documents:
        return jsonify({'error': 'No documents provided'}), 400

    try:
        processed_documents = []
        for doc in documents:
            content = base64.b64decode(doc['content'])
            file_type = doc['type'].lower()

            if file_type == 'pdf':
                text = extract_text(io.BytesIO(content))
            elif file_type in ['doc', 'docx']:
                doc = Document(io.BytesIO(content))
                text = "\n".join([para.text for para in doc.paragraphs])
            elif file_type in ['xls', 'xlsx']:
                wb = load_workbook(io.BytesIO(content))
                text = ""
                for sheet in wb:
                    for row in sheet.iter_rows(values_only=True):
                        text += " | ".join([str(cell) for cell in row if cell]) + "\n"
            elif file_type in ['ppt', 'pptx']:
                prs = Presentation(io.BytesIO(content))
                text = ""
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, 'text'):
                            text += shape.text + "\n"
            elif file_type in ['jpg', 'jpeg', 'png', 'gif', 'bmp']:
                image = Image.open(io.BytesIO(content))
                text = pytesseract.image_to_string(image)
            else:
                text = content.decode('utf-8', errors='ignore')

            processed_documents.append({
                'name': doc['name'],
                'content': text,
                'type': file_type
            })

        summaries = summarize_documents(summarization_model, processed_documents, merge_summaries, summary_depth, language)
        
        if enable_visual_enhancements:
            enriched_summaries = [enrich_summary_with_visuals(summary) for summary in summaries]
        else:
            enriched_summaries = summaries

        response = {
            'user_id': user_id,
            'summaries': enriched_summaries
        }
        return jsonify(response), 200
    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/feedback', methods=['POST'])
def feedback():
    data = request.json
    summary_id = data.get('summary_id')
    user_id = data.get('user_id')
    feedback_text = data.get('feedback')

    if not all([summary_id, user_id, feedback_text]):
        return jsonify({'error': 'Missing required data'}), 400

    try:
        # Here you would typically update the model based on the feedback
        # For now, we'll just acknowledge the feedback
        return jsonify({'message': 'Feedback received successfully'}), 200
    except Exception as e:
        return jsonify({'error': str(e)}), 500

if __name__ == '__main__':
    app.run(debug=True)